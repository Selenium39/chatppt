import io
import json
import os
import re
import uuid

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from urllib.parse import quote_plus
from dotenv import load_dotenv

dir_path = 'static/presentations'

load_dotenv()
API_KEY = os.getenv('PEXELS_API_KEY')

def parse_response(response):
    match = re.search(r'\[.*\]', response, re.DOTALL)
    if match is None:
        print("No JSON array found in response.")
        return None
    json_str = match.group(0)
    try:
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        print(f"An error occurred while parsing the JSON array: {e}")
        return None



def search_pexels_images(keyword):
    keyword = keyword[0].lower() if keyword else ""
    query = quote_plus(keyword)
    print("Query:", query) # Debug
    PEXELS_API_URL = f'https://api.pexels.com/v1/search?query={query}&per_page=1&locale=zh-CN'
    print("URL:", PEXELS_API_URL) # Debug
    headers = {
        'Authorization': API_KEY
    }
    response = requests.get(PEXELS_API_URL, headers=headers)
    print("Response Status Code:", response.status_code) # Debug
    print("Response Content:", response.text) # Debug
    data = json.loads(response.text)
    if 'photos' in data:
        if len(data['photos']) > 0:
            return data['photos'][0]['src']['medium']
    return None


def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(slides_content, template_choice, presentation_title, presenter_name, insert_image):
    template_path = os.path.join(dir_path, f"{template_choice}.pptx")

    prs = Presentation(template_path)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title

    #add subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = f"由 {presenter_name} 呈现"

    if template_choice == 'dark_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color

    elif template_choice == 'bright_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Arial'
                run.font.color.rgb = RGBColor(255, 20, 147)  # RGB for deep pink color

    # add content slides
    for slide_content in slides_content:
        print(slide_content)
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = slide_content['content']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for white color

        if insert_image:
            # fetch image URL from Pixabay based on the slide's title
            image_url = search_pexels_images(slide_content['keyword'])
            print("Image URL:", image_url) #debug
            if image_url is not None:
                # download the image
                image_data = requests.get(image_url).content
                # load image into BytesIO object
                image_stream = io.BytesIO(image_data)
                # add the image at the specified position
                slide_width = Inches(20)
                slide_height = Inches(15)

                image_width = Inches(8)  # width of image
                image_height = Inches(5)  # height of image

                left = slide_width - image_width  # calculate left position
                top = slide_height - image_height - Inches(4)  # calculate top position

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    unique_filename = f"{uuid.uuid4()}.pptx"
    # Save the presentation
    prs.save(os.path.join('generated', unique_filename))
    return unique_filename
